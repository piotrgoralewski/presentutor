from flask import Flask, render_template, request, send_file
from pptx import Presentation


app = Flask(__name__)
slides_content = {}


@app.route('/')
def hello_world():
    """
    Loads the main page
    :return:
    """
    return render_template("index.html")


@app.route('/set_content', methods=["POST"])
def set_content():
    """
    Sets the slides_content variable (this has to be JSON in the format as described by the chatbot). This endpoint is
    called by the event when Watson Assistant set up variable "slideshow_json".

    "JSON format, where each field name is a title for a slide, and the value is either a list of bullet points, or
    just a string, for a slide with text only."
    :return:
    """
    global slides_content  # todo: This should be object variable instead of global...
    slides_content = request.get_json()
    return "ok"


@app.route("/make_presentation")
def make_presentation():
    """
    Endpoint called to create and return (as download file) PowerPoint presentation.
    :return:
    """
    global slides_content
    prs = Presentation()
    for title, content in slides_content.items():
        if isinstance(content, list):  # If content is a list, create a bullet point slide
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            title_placeholder = slide.shapes.title
            content_placeholder = slide.placeholders[1]

            title_placeholder.text = title
            for bullet_point in content:
                p = content_placeholder.text_frame.add_paragraph()
                p.text = bullet_point
                p.level = 0  # You can increase the level for sub-bullets
        else:  # If content is a string, create a simple text slide
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            title_placeholder = slide.shapes.title
            content_placeholder = slide.placeholders[1]

            title_placeholder.text = title
            content_placeholder.text = content

    # Save the presentation
    prs.save('Presentation.pptx')
    return send_file("Presentation.pptx", download_name='Presentation.pptx', as_attachment=True)


if __name__ == '__main__':
    app.run()
